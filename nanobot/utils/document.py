"""Document text extraction utilities for nanobot."""

import mimetypes
from pathlib import Path

from loguru import logger

from nanobot.utils.helpers import detect_image_mime


# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS: set[str] = {
    # Document formats
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    # Text formats
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    # Image formats (for future OCR support)
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
}

def extract_text(path: Path,max_size:int) -> str | None:
    """Extract text from a file.

    Args:
        path: Path to the file.
        max_size: Maximum size in file content bytes.

    Returns:
        Extracted text as string, None for unsupported types,
        or error string for failures.
    """
    if not isinstance(path, Path):
        path = Path(path)

    if not path.exists():
        return f"[error: file not found: {path}]"

    ext = path.suffix.lower()

    # Document formats -- each branch lazily imports its parser so that
    # startup does not pay the ~25 MB cost of loading openpyxl /
    # python-docx / python-pptx / pypdf up front (see issue #3422).
    if ext == ".pdf":
        return _extract_pdf(path,max_size)
    elif ext == ".docx":
        return _extract_docx(path,max_size)
    elif ext == ".xlsx":
        return _extract_xlsx(path,max_size)
    elif ext == ".pptx":
        return _extract_pptx(path,max_size)
    elif _is_text_extension(ext):
        return _extract_text_file(path,max_size)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        # Image files - for future OCR support
        return f"[image: {path.name}]"
    else:
        # Unsupported extension
        return None


def _extract_pdf(path: Path,max_size:int) -> str:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return "[error: pypdf not installed]"
    try:
        reader = PdfReader(path)
        pages: list[str] = []
        current_total_size = 0

        for i, page in enumerate(reader.pages, 1):
            # extract one page
            page_text = page.extract_text() or ""
            page_header = f"--- Page {i} ---\n"
            combined_page = page_header + page_text + "\n\n"

            # check current size
            if current_total_size + len(combined_page) >= max_size:
                remaining_space = max_size - current_total_size
                if remaining_space > len(page_header):
                    pages.append(combined_page[:remaining_space])
                break

            # calculate current size
            pages.append(combined_page)
            current_total_size += len(combined_page)

        return "".join(pages).strip()
    except Exception as e:
        logger.exception("Failed to extract PDF {}", path)
        return f"[error: failed to extract PDF: {e!s}]"


def _extract_docx(path: Path,max_size:int) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return "[error: python-docx not installed]"
    try:
        doc = DocxDocument(str(path))
        paragraphs: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        current_size = 0

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue

            # Check if adding this paragraph exceeds the limit
            if current_size + len(text) > max_size:
                remaining_space = max_size - current_size
                if remaining_space > 0:
                    paragraphs.append(text[:remaining_space])
                break  # Stop processing immediately to save memory

            paragraphs.append(text)
            current_size += len(text) + 2  # +2 accounts for the "\n\n" joiner

        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.exception("Failed to extract DOCX {}", path)
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path,max_size:int) -> str:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[error: openpyxl not installed]"
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        current_total_size=0
        try:
            sheets: list[str] = []
            for sheet_name in wb.sheetnames:
                if current_total_size >= max_size:
                    break

                ws = wb[sheet_name]
                header = f"--- Sheet: {sheet_name} ---\n"
                sheets.append(header)
                current_total_size += len(header)

                #
                for row in ws.iter_rows(values_only=True):
                    # Join cell values with tabs, skipping empty cells
                    raw_line = "\t".join(str(c) if c is not None else "" for c in row).strip()
                    if not raw_line:
                        continue

                    line = raw_line + "\n"
                    length = len(line)

                    if current_total_size + length > max_size:
                        #
                        sheets.append(line[:max_size - current_total_size])
                        current_total_size = max_size
                        break

                    sheets.append(line)
                    current_total_size += length

                sheets.append("\n")
                current_total_size += 1

            return "".join(sheets)
        finally:
            wb.close()
    except Exception as e:
        logger.exception("Failed to extract XLSX {}", path)
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path, max_size: int) -> str:
    """Extract text from PPTX while strictly limiting memory usage."""
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        return "[error: python-pptx not installed]"

    try:
        prs = PptxPresentation(str(path))
        extracted_content:list[str] = []
        current_size = 0

        for i, slide in enumerate(prs.slides, 1):
            slide_header = f"--- Slide {i} ---\n"
            slide_text_list: list[str] = []

            # range ppt
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text_list, max_size - current_size - len(slide_header))

            if slide_text_list:
                combined_slide = slide_header + "\n".join(slide_text_list)
                extracted_content.append(combined_slide)
                current_size += len(combined_slide) + 2  # joiner "\n\n so plus 2"

            if current_size >= max_size:
                break

        return "\n\n".join(extracted_content)[:max_size]
    except Exception as e:
        logger.exception("Failed to extract PPTX {}", path)
        return f"[error: failed to extract PPTX: {e!s}]"


def _collect_pptx_shape_text(shape, out: list[str], remaining_quota: int) -> None:
    """Collect text with a safety quota check."""
    current_out_size = sum(len(s) for s in out)
    if current_out_size >= remaining_quota:
        return

    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out, remaining_quota)
        return

    # handle table
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                if (sum(len(s) for s in out) + len(line)) > remaining_quota:
                    break
                out.append(line)
        return

    # handle text file
    text = getattr(shape, "text", "").strip()
    if text:
        if (sum(len(s) for s in out) + len(text)) > remaining_quota:
            # even though a shape, cut here if it's too long
            allowable = remaining_quota - sum(len(s) for s in out)
            if allowable > 0:
                out.append(text[:allowable])
        else:
            out.append(text)


def _extract_text_file(path: Path, max_size: int) -> str:
    """Extract text from a plain text file with memory-safe streaming."""
    # Attempt UTF-8 first, fallback to latin-1 if decoding fails
    for encoding in ["utf-8", "latin-1"]:
        try:
            with path.open("r", encoding=encoding, errors="replace") as f:
                # Read only up to max_size characters
                # This prevents loading the entire file into RAM
                return f.read(max_size)
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logger.exception("Failed to read text file {}", path)
            return f"[error: failed to read file: {e!s}]"

    return "[error: could not decode file with supported encodings]"


def _is_text_extension(ext: str) -> bool:
    """Check if extension is a text format."""
    return ext in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
    }


# ---------------------------------------------------------------------------
# High-level helper: split media into images + extracted document text
# ---------------------------------------------------------------------------

_MAX_EXTRACT_FILE_SIZE = 50 * 1024 * 1024  # 50 MB


def extract_documents(
    text: str,
    media_paths: list[str],
    *,
    max_file_size: int = _MAX_EXTRACT_FILE_SIZE,
) -> tuple[str, list[str]]:
    """Separate images from documents in *media_paths*.

    Documents (PDF, DOCX, XLSX, PPTX, plain-text, …) have their text
    extracted and appended to *text*.  Only image paths are kept in the
    returned list so that downstream layers only need to handle vision
    blocks.

    Files larger than *max_file_size* bytes are skipped with a warning
    to avoid unbounded memory / CPU usage.
    """
    image_paths: list[str] = []
    doc_texts: list[str] = []

    for path_str in media_paths:
        p = Path(path_str)
        if not p.is_file():
            continue

        try:
            size = p.stat().st_size
        except OSError:
            continue
        if size > max_file_size:
            logger.warning(
                "Skipping oversized file for extraction: {} ({:.1f} MB > {} MB limit)",
                p.name, size / (1024 * 1024), max_file_size // (1024 * 1024),
            )
            continue

        with open(p, "rb") as f:
            header = f.read(16)
        mime = detect_image_mime(header) or mimetypes.guess_type(path_str)[0]
        if mime and mime.startswith("image/"):
            image_paths.append(path_str)
        else:
            extracted = extract_text(p)
            if extracted and not extracted.startswith("[error:"):
                doc_texts.append(f"[File: {p.name}]\n{extracted}")

    if doc_texts:
        text = text + "\n\n" + "\n\n".join(doc_texts)

    return text, image_paths
